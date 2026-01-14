#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HTML Presentation to PowerPoint Converter
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re

def clean_text(text):
    """Remove extra whitespace and clean text"""
    if not text:
        return ""
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def add_text_frame(slide, left, top, width, height, text, font_size=14, bold=False, color=RGBColor(0, 0, 0)):
    """Add a text box to the slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color

    return textbox

def convert_html_to_pptx(html_file, output_file):
    """Convert HTML presentation to PowerPoint"""

    # Read HTML file
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()

    # Parse HTML
    soup = BeautifulSoup(html_content, 'html.parser')

    # Create PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 aspect ratio
    prs.slide_height = Inches(5.625)

    # Find all slides (sections)
    slides = soup.find_all('section', recursive=True)

    print(f"Found {len(slides)} slides")

    for idx, section in enumerate(slides, 1):
        print(f"Processing slide {idx}...")

        # Create blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Position variables
        left_margin = Inches(0.5)
        right_margin = Inches(0.5)
        top = Inches(0.5)
        width = prs.slide_width - left_margin - right_margin
        line_height = Inches(0.3)

        # Extract and add title (h1, h2)
        title = section.find(['h1', 'h2'])
        if title:
            title_text = clean_text(title.get_text())
            font_size = 32 if title.name == 'h1' else 28
            add_text_frame(slide, left_margin, top, width, Inches(0.8),
                          title_text, font_size=font_size, bold=True,
                          color=RGBColor(59, 130, 246))
            top += Inches(1.0)

        # Extract content
        content_parts = []

        # Get all child elements
        for element in section.find_all(['h3', 'p', 'ul', 'li', 'div', 'table']):
            # Skip if this element is inside another element we'll process
            if element.find_parent(['ul']):
                continue

            if element.name == 'h3':
                text = clean_text(element.get_text())
                if text:
                    content_parts.append(('h3', text))

            elif element.name == 'p':
                # Skip if parent is li (will be handled with li)
                if element.parent.name not in ['li', 'td', 'th']:
                    text = clean_text(element.get_text())
                    if text and len(text) > 3:  # Skip very short texts
                        content_parts.append(('p', text))

            elif element.name == 'ul':
                # Process list
                for li in element.find_all('li', recursive=False):
                    text = clean_text(li.get_text())
                    if text:
                        content_parts.append(('li', text))

            elif element.name == 'div' and 'highlight-box' in element.get('class', []):
                text = clean_text(element.get_text())
                if text:
                    content_parts.append(('highlight', text))

            elif element.name == 'table':
                # Add a note that there's a table
                content_parts.append(('table', '[表: 詳細はHTMLを参照]'))
                # Try to extract first row as summary
                thead = element.find('thead') or element.find('tr')
                if thead:
                    headers = [clean_text(th.get_text()) for th in thead.find_all(['th', 'td'])]
                    if headers:
                        content_parts.append(('p', ' | '.join(headers)))

        # Add content to slide
        current_top = top
        max_height = prs.slide_height - Inches(0.5)

        for content_type, text in content_parts:
            # Check if we have space
            if current_top >= max_height:
                break

            if content_type == 'h3':
                add_text_frame(slide, left_margin, current_top, width, Inches(0.4),
                              text, font_size=20, bold=True, color=RGBColor(59, 130, 246))
                current_top += Inches(0.5)

            elif content_type == 'p':
                # Estimate height based on text length
                estimated_lines = max(1, len(text) // 80)
                text_height = Inches(0.2 + estimated_lines * 0.15)
                add_text_frame(slide, left_margin, current_top, width, text_height,
                              text, font_size=14, color=RGBColor(100, 116, 139))
                current_top += text_height + Inches(0.1)

            elif content_type == 'li':
                # Add bullet point
                text = "• " + text
                estimated_lines = max(1, len(text) // 80)
                text_height = Inches(0.2 + estimated_lines * 0.15)
                add_text_frame(slide, left_margin + Inches(0.2), current_top,
                              width - Inches(0.2), text_height,
                              text, font_size=14, color=RGBColor(100, 116, 139))
                current_top += text_height + Inches(0.05)

            elif content_type == 'highlight':
                # Add highlighted text
                estimated_lines = max(1, len(text) // 80)
                text_height = Inches(0.3 + estimated_lines * 0.15)
                textbox = add_text_frame(slide, left_margin, current_top, width, text_height,
                                        text, font_size=16, bold=True, color=RGBColor(30, 41, 59))
                # Add background color
                fill = textbox.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(239, 246, 255)
                current_top += text_height + Inches(0.15)

            elif content_type == 'table':
                add_text_frame(slide, left_margin, current_top, width, Inches(0.3),
                              text, font_size=12, color=RGBColor(156, 163, 175))
                current_top += Inches(0.4)

    # Save presentation
    prs.save(output_file)
    print(f"\nPowerPoint file saved: {output_file}")
    print(f"Total slides created: {len(prs.slides)}")

if __name__ == "__main__":
    html_file = "/home/user/Review-Predictor/発表用HTML/finalapp/presentation.html"
    output_file = "/home/user/Review-Predictor/発表用HTML/finalapp/Review_Predictor_Presentation.pptx"

    convert_html_to_pptx(html_file, output_file)
